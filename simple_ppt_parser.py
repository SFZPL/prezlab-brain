import io
import os
import subprocess
import uuid
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = Flask(__name__)
CORS(app)

# Allow large uploads (configurable via env, default 600MB)
max_mb = int(os.getenv('PARSER_MAX_MB', '600'))
app.config['MAX_CONTENT_LENGTH'] = max_mb * 1024 * 1024

# Thumbnails root directory
THUMBS_ROOT = os.path.abspath(os.getenv('PARSER_THUMBS_DIR', os.path.join(os.path.dirname(__file__), 'thumb_cache')))
os.makedirs(THUMBS_ROOT, exist_ok=True)


def _get_alt_text(shape):
    try:
        if hasattr(shape, "_element") and hasattr(shape._element, "nvPicPr"):
            return shape._element.nvPicPr.cNvPr.get('descr', '') or ''
        if hasattr(shape, "_element") and hasattr(shape._element, "nvSpPr"):
            return shape._element.nvSpPr.cNvPr.get('descr', '') or ''
    except Exception:
        return ''
    return ''


def _analyze_shape(shape):
    counters = {
        'images': 0,
        'charts': 0,
        'tables': 0,
        'text_boxes': 0,
        'shapes_count': 1,
        'alt_texts': []
    }
    text_lines = []

    try:
        stype = getattr(shape, 'shape_type', None)

        if stype == MSO_SHAPE_TYPE.GROUP and hasattr(shape, 'shapes'):
            for child in shape.shapes:
                c_counters, c_text = _analyze_shape(child)
                for k in counters:
                    if k == 'alt_texts':
                        counters[k].extend(c_counters[k])
                    else:
                        counters[k] += c_counters[k]
                text_lines.extend(c_text)
            return counters, text_lines

        if stype == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, 'image'):
            counters['images'] += 1
            alt_text = _get_alt_text(shape)
            if alt_text:
                counters['alt_texts'].append(alt_text)

        if hasattr(shape, 'has_chart') and shape.has_chart:
            counters['charts'] += 1

        if hasattr(shape, 'has_table') and shape.has_table:
            counters['tables'] += 1

        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            counters['text_boxes'] += 1
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs).strip()
                if line:
                    text_lines.append(line)
    except Exception:
        pass

    return counters, text_lines


def extract_pptx_data(file_stream):
    prs = Presentation(file_stream)

    slides_data = []
    all_text_lines = []

    for idx, slide in enumerate(prs.slides, start=1):
        total_counters = {
            'images': 0,
            'charts': 0,
            'tables': 0,
            'text_boxes': 0,
            'shapes_count': 0,
            'alt_texts': []
        }
        text_lines = []

        for shape in slide.shapes:
            c_counters, c_text = _analyze_shape(shape)
            for k in total_counters:
                if k == 'alt_texts':
                    total_counters[k].extend(c_counters[k])
                else:
                    total_counters[k] += c_counters[k]
            text_lines.extend(c_text)

        slide_text = "\n".join(text_lines)
        all_text_lines.extend(text_lines)

        text_chars = len(slide_text)
        visual_score = (total_counters['images'] * 3) + (total_counters['charts'] * 4) + (total_counters['tables'] * 2)
        density = 'visual-heavy' if visual_score >= 6 and text_chars < 800 else 'balanced' if visual_score >= 3 else 'text-heavy'

        slides_data.append({
            "slide_number": idx,
            "text_content": slide_text,
            "layout_type": getattr(slide.slide_layout, 'name', 'unknown') or 'unknown',
            "images": total_counters['images'],
            "charts": total_counters['charts'],
            "tables": total_counters['tables'],
            "text_boxes": total_counters['text_boxes'],
            "shapes_count": total_counters['shapes_count'],
            "alt_texts": total_counters['alt_texts'][:10],
            "notes": getattr(slide.notes_slide.notes_text_frame, 'text', '').strip() if hasattr(slide, 'notes_slide') and slide.notes_slide and hasattr(slide.notes_slide, 'notes_text_frame') and slide.notes_slide.notes_text_frame else "",
            "visual_density": density
        })

    text_content = "\n".join(all_text_lines)

    return {
        "slide_count": len(slides_data),
        "text_content": text_content,
        "slides": slides_data,
        "metadata": {
            "parseMethod": "python-pptx",
            "visual_summary": {
                "total_images": sum(s['images'] for s in slides_data),
                "total_charts": sum(s['charts'] for s in slides_data),
                "total_tables": sum(s['tables'] for s in slides_data)
            }
        }
    }


def _render_with_powerpoint_com(input_path: str, out_dir: str) -> bool:
    """Render thumbnails using Microsoft PowerPoint COM automation (Windows only)."""
    try:
        if os.name != 'nt':
            return False
        if os.getenv('PARSER_USE_PPT_COM', '0') != '1':
            return False

        # Lazy import so non-Windows envs don't fail
        try:
            from win32com.client import Dispatch
        except Exception:
            return False

        # Width/height optional overrides
        try:
            width = int(os.getenv('PARSER_PNG_WIDTH', '1920'))
            height = int(os.getenv('PARSER_PNG_HEIGHT', '1080'))
        except Exception:
            width, height = 1920, 1080

        ppt = Dispatch('PowerPoint.Application')
        ppt.Visible = 0
        try:
            pres = ppt.Presentations.Open(input_path, WithWindow=False, ReadOnly=True)
            try:
                # Export each slide as PNG
                pres.Export(out_dir, 'PNG', width, height)
            finally:
                pres.Close()
        finally:
            ppt.Quit()

        # Ensure some PNGs exist
        pngs = [fn for fn in os.listdir(out_dir) if fn.lower().endswith('.png')]
        pngs.sort()
        return len(pngs) > 0
    except Exception:
        return False


def _render_thumbnails_to_disk(file_bytes, original_filename):
    """Generate PNG thumbnails via COM (if enabled) or LibreOffice."""
    try:
        job_id = str(uuid.uuid4())
        job_dir = os.path.join(THUMBS_ROOT, job_id)
        os.makedirs(job_dir, exist_ok=True)
        input_path = os.path.join(job_dir, secure_filename(original_filename))
        with open(input_path, 'wb') as f:
            f.write(file_bytes)

        # 1) Try high-fidelity PowerPoint COM rendering if enabled
        if _render_with_powerpoint_com(input_path, job_dir):
            pngs = [fn for fn in os.listdir(job_dir) if fn.lower().endswith('.png')]
            pngs.sort()
            return job_id, pngs

        # 2) Fallback to LibreOffice
        soffice_path = os.getenv('SOFFICE_PATH', 'soffice')
        cmd = [soffice_path, '--headless', '--convert-to', 'png:impress_png_Export', '--outdir', job_dir, input_path]
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=int(os.getenv('PARSER_RENDER_TIMEOUT', '180')))
        if result.returncode != 0:
            return None, []

        pngs = [fn for fn in os.listdir(job_dir) if fn.lower().endswith('.png')]
        pngs.sort()
        return job_id, pngs
    except Exception:
        return None, []


@app.route('/thumbnails/<path:subpath>')
def thumbnails(subpath):
    full_path = os.path.join(THUMBS_ROOT, subpath)
    directory = os.path.dirname(full_path)
    filename = os.path.basename(full_path)
    return send_from_directory(directory, filename)


@app.route('/health', methods=['GET'])
def health():
    return jsonify({"status": "ok", "max_mb": max_mb}), 200


@app.route('/parse-pptx', methods=['POST'])
def parse_pptx():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    filename = secure_filename(file.filename)
    ext = (os.path.splitext(filename)[1] or '').lower()
    if ext not in ['.pptx', '.ppt']:
        return jsonify({"error": "Unsupported file type. Please upload a .pptx or .ppt file."}), 400

    try:
        file_bytes = file.read()
        stream = io.BytesIO(file_bytes)

        data = extract_pptx_data(stream)

        # Optional rendering of thumbnails
        should_render = os.getenv('PARSER_RENDER_THUMBS', '1') != '0'
        job_id, pngs = (None, [])
        if should_render:
            job_id, pngs = _render_thumbnails_to_disk(file_bytes, filename)

        if job_id and pngs:
            base = request.host_url.rstrip('/')
            # Attach per-slide thumbnail_url (best-effort order match)
            for idx, slide in enumerate(data['slides']):
                if idx < len(pngs):
                    slide['thumbnail_url'] = f"{base}/thumbnails/{job_id}/{pngs[idx]}"

        if not data.get('text_content') and not data.get('slides'):
            return jsonify({"error": "No extractable content found in the presentation."}), 422

        return jsonify(data), 200
    except Exception as e:
        return jsonify({"error": f"Parsing failed: {str(e)}"}), 500


@app.route('/analyze-slide', methods=['POST'])
def analyze_slide():
    try:
        data = request.get_json(force=True)
        slide = data.get('slide_data', {})
        presentation_context = data.get('presentation_context', {})

        text = (slide.get('text_content') or '').strip()
        word_count = len(text.split()) if text else 0
        images = int(slide.get('images') or 0)
        charts = int(slide.get('charts') or 0)
        tables = int(slide.get('tables') or 0)

        length_label = 'appropriate'
        if word_count > 90:
            length_label = 'too long'
        elif word_count < 10 and (images + charts + tables) == 0:
            length_label = 'too short'

        clarity = 'clear' if word_count <= 50 else 'mixed'
        organization = 'structured' if '\n' in text or '-' in text or '•' in text else 'block text'
        effectiveness = 'high' if images + charts + tables > 0 or word_count <= 60 else 'medium'
        priority = 'important' if length_label != 'appropriate' else 'normal'

        response = {
            "slideOverview": {
                "slideNumber": slide.get('slide_number') or 1,
                "contentSummary": text[:180] + ('...' if len(text) > 180 else ''),
                "slidePurpose": "content" if word_count > 0 else "visual",
                "effectiveness": effectiveness,
                "priority": priority
            },
            "contentAnalysis": {
                "textContent": {
                    "clarity": clarity,
                    "organization": organization,
                    "length": length_label,
                    "keyMessages": [text[:90]] if text else [],
                    "improvements": [
                        "Reduce text density" if length_label == 'too long' else "",
                        "Add structure with bullets" if organization == 'block text' else "",
                        "Add a supporting visual" if (images + charts + tables) == 0 else ""
                    ]
                },
                "visualElements": {
                    "images": {"count": images, "relevance": "unknown"},
                    "charts": {"count": charts, "effectiveness": "unknown"},
                    "tables": {"count": tables, "readability": "unknown"}
                }
            },
            "designRecommendations": {
                "layout": {
                    "currentLayout": slide.get('layout_type') or 'unknown',
                    "effectiveness": "unknown",
                    "recommendedLayout": "title + content" if word_count > 0 else "full-bleed visual",
                    "specificChanges": [
                        "Use bullets for lists" if organization == 'block text' else "",
                        "Increase whitespace around content",
                        "Limit to one idea per slide"
                    ],
                    "visualHierarchy": "headline > key point > support"
                },
                "colorScheme": {
                    "currentColors": slide.get('colors') or [],
                    "effectiveness": "unknown",
                    "recommendedPalette": ["#1F2937", "#3B82F6", "#F3F4F6"],
                    "colorPsychology": "neutral with highlight",
                    "accessibility": "ensure 4.5:1 contrast"
                },
                "typography": {
                    "currentFonts": slide.get('fonts') or [],
                    "readability": "unknown",
                    "recommendedFonts": ["Inter", "Source Sans Pro"],
                    "sizing": "Headline 28-36, body 16-18",
                    "hierarchy": "headline > bullets > notes"
                }
            },
            "quickWins": [
                "Convert paragraphs to 3-5 bullets",
                "Highlight 3-5 keywords",
                "Add a relevant visual"
            ],
            "priorityFixes": [
                "Reduce word count" if length_label == 'too long' else "Ensure one message per slide"
            ]
        }

        response["contentAnalysis"]["textContent"]["improvements"] = [
            i for i in response["contentAnalysis"]["textContent"]["improvements"] if i
        ]

        return jsonify(response), 200

    except Exception as e:
        return jsonify({"error": f"Slide analysis failed: {str(e)}"}), 500


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=False)
