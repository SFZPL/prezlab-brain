from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import os
import json
import tempfile
import shutil
from werkzeug.utils import secure_filename
import openai
from pptx import Presentation
import mammoth
import PyPDF2
import io
import base64
from datetime import datetime
import logging
from dotenv import load_dotenv
from typing import Dict, List, Any, Optional
from concurrent.futures import ThreadPoolExecutor
from enum import Enum
from functools import wraps
import time
import gc
import psutil
import hashlib
import pickle

# Enhanced error handling and performance improvements
import traceback
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__, static_folder='build', static_url_path='')
CORS(app)

# Production configuration
if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)

# Enhanced configuration
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024 * 3  # 600MB max file size
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0

# Create upload folder if it doesn't exist
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Error types
class ErrorType(Enum):
    FILE_TOO_LARGE = "FILE_TOO_LARGE"
    PARSING_FAILED = "PARSING_FAILED"
    API_ERROR = "API_ERROR"
    TIMEOUT = "TIMEOUT"
    INVALID_FORMAT = "INVALID_FORMAT"
    COLOR_EXTRACTION_FAILED = "COLOR_EXTRACTION_FAILED"
    MEMORY_ERROR = "MEMORY_ERROR"

class PresentationAnalyzerError(Exception):
    """Custom exception for presentation analyzer"""
    def __init__(self, error_type: ErrorType, message: str, details: Optional[Dict] = None):
        self.error_type = error_type
        self.message = message
        self.details = details or {}
        super().__init__(self.message)

# Rate limiting
class RateLimiter:
    """Rate limiting for API calls"""
    def __init__(self, max_calls: int = 10, time_window: int = 60):
        self.max_calls = max_calls
        self.time_window = time_window
        self.calls = []
    
    def can_make_call(self) -> bool:
        now = time.time()
        # Remove old calls outside time window
        self.calls = [call_time for call_time in self.calls 
                     if now - call_time < self.time_window]
        
        if len(self.calls) < self.max_calls:
            self.calls.append(now)
            return True
        return False

# Memory management
class MemoryManager:
    """Monitor and manage memory usage"""
    
    @staticmethod
    def check_memory_usage():
        """Check current memory usage"""
        try:
            process = psutil.Process()
            memory_info = process.memory_info()
            return {
                'rss': memory_info.rss / 1024 / 1024,  # MB
                'percent': process.memory_percent()
            }
        except Exception as e:
            logger.warning(f"Could not check memory usage: {e}")
            return {'rss': 0, 'percent': 0}
    
    @staticmethod
    def cleanup_if_needed():
        """Force garbage collection if memory usage is high"""
        try:
            memory = MemoryManager.check_memory_usage()
            if memory['percent'] > 80:
                gc.collect()
                logger.info(f"Forced garbage collection. Memory usage: {memory['percent']:.2f}%")
        except Exception as e:
            logger.warning(f"Memory cleanup failed: {e}")
    
    @staticmethod
    def process_with_memory_management(func):
        """Decorator to manage memory during processing"""
        @wraps(func)
        def wrapper(*args, **kwargs):
            try:
                # Check memory before
                MemoryManager.cleanup_if_needed()
                
                # Process
                result = func(*args, **kwargs)
                
                # Cleanup after
                MemoryManager.cleanup_if_needed()
                
                return result
            except MemoryError:
                gc.collect()
                raise PresentationAnalyzerError(
                    ErrorType.MEMORY_ERROR,
                    "Insufficient memory to process this file",
                    {"suggestion": "Try processing a smaller file or fewer slides at once"}
                )
        return wrapper

# File validation
class FileValidator:
    """Validate files before processing"""
    
    MAX_FILE_SIZE = 600 * 1024 * 1024  # 600MB
    ALLOWED_EXTENSIONS = {'ppt', 'pptx', 'doc', 'docx', 'pdf', 'txt'}
    
    @classmethod
    def validate_file(cls, file) -> Dict[str, Any]:
        """Comprehensive file validation"""
        errors = []
        warnings = []
        
        # Check file size
        file.seek(0, 2)  # Seek to end
        file_size = file.tell()
        file.seek(0)  # Reset to beginning
        
        if file_size > cls.MAX_FILE_SIZE:
            errors.append({
                'type': ErrorType.FILE_TOO_LARGE,
                'message': f'File size ({file_size / 1024 / 1024:.2f}MB) exceeds maximum allowed size ({cls.MAX_FILE_SIZE / 1024 / 1024}MB)'
            })
        elif file_size > 100 * 1024 * 1024:  # Warning for files > 100MB
            warnings.append({
                'type': 'LARGE_FILE',
                'message': f'Large file detected ({file_size / 1024 / 1024:.2f}MB). Processing may take longer.'
            })
        
        # Check file extension
        filename = file.filename
        if not filename:
            errors.append({
                'type': ErrorType.INVALID_FORMAT,
                'message': 'No filename provided'
            })
        else:
            ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
            if ext not in cls.ALLOWED_EXTENSIONS:
                errors.append({
                    'type': ErrorType.INVALID_FORMAT,
                    'message': f'File type .{ext} is not supported. Supported types: {", ".join(cls.ALLOWED_EXTENSIONS)}'
                })
        
        # Check file content (magic bytes)
        file_header = file.read(8)
        file.seek(0)
        
        # PowerPoint magic bytes
        if ext in ['ppt', 'pptx']:
            if not (file_header.startswith(b'\xd0\xcf\x11\xe0') or  # PPT
                   file_header.startswith(b'PK')):  # PPTX (ZIP)
                warnings.append({
                    'type': 'CORRUPTED_FILE',
                    'message': 'File may be corrupted or not a valid PowerPoint file'
                })
        
        return {
            'valid': len(errors) == 0,
            'errors': errors,
            'warnings': warnings,
            'file_info': {
                'size': file_size,
                'extension': ext,
                'filename': filename
            }
        }

# Error handling decorator
class ErrorHandler:
    """Centralized error handling with recovery strategies"""
    
    @staticmethod
    def handle_parsing_error(func):
        """Decorator for handling parsing errors with fallback strategies"""
        @wraps(func)
        def wrapper(*args, **kwargs):
            attempts = 3
            fallback_strategies = [
                lambda: func(*args, **kwargs),  # Original attempt
                lambda: ErrorHandler._parse_with_reduced_features(*args, **kwargs),  # Reduced features
                lambda: ErrorHandler._parse_text_only(*args, **kwargs)  # Text only
            ]
            
            for attempt, strategy in enumerate(fallback_strategies):
                try:
                    return strategy()
                except Exception as e:
                    logger.warning(f"Attempt {attempt + 1} failed: {str(e)}")
                    if attempt == len(fallback_strategies) - 1:
                        # All strategies failed
                        raise PresentationAnalyzerError(
                            ErrorType.PARSING_FAILED,
                            "Unable to parse presentation after multiple attempts",
                            {"original_error": str(e), "attempts": attempt + 1}
                        )
            
        return wrapper
    
    @staticmethod
    def _parse_with_reduced_features(file_path, *args, **kwargs):
        """Parse with reduced features (skip colors, fonts, etc.)"""
        parser = EnhancedFileParser()
        parser.skip_formatting = True  # Add this flag to your parser
        return parser.parse_powerpoint_chunked(file_path)
    
    @staticmethod
    def _parse_text_only(file_path, *args, **kwargs):
        """Extract text only as last resort"""
        from pptx import Presentation
        prs = Presentation(file_path)
        text_content = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text.append(shape.text)
            text_content.append('\n'.join(slide_text))
        
        return {
            'slide_count': len(prs.slides),
            'text_content': '\n\n'.join(text_content),
            'metadata': {'parse_method': 'text_only_fallback'},
            'slides': []
        }

# Initialize rate limiter
rate_limiter = RateLimiter(max_calls=20, time_window=60)

# Error handling functions
def get_error_status_code(error_type: ErrorType) -> int:
    """Get appropriate HTTP status code for error type"""
    status_codes = {
        ErrorType.FILE_TOO_LARGE: 413,
        ErrorType.INVALID_FORMAT: 415,
        ErrorType.TIMEOUT: 504,
        ErrorType.PARSING_FAILED: 422,
        ErrorType.API_ERROR: 502,
        ErrorType.COLOR_EXTRACTION_FAILED: 422,
        ErrorType.MEMORY_ERROR: 507
    }
    return status_codes.get(error_type, 500)

def get_error_suggestions(error_type: ErrorType) -> list:
    """Provide helpful suggestions based on error type"""
    suggestions = {
        ErrorType.FILE_TOO_LARGE: [
            "Split your presentation into smaller files",
            "Remove high-resolution images",
            "Export as PDF and try again"
        ],
        ErrorType.PARSING_FAILED: [
            "Ensure the file is not corrupted",
            "Try saving in a different format (PPTX instead of PPT)",
            "Remove complex animations or transitions",
            "Export as PDF for basic analysis"
        ],
        ErrorType.COLOR_EXTRACTION_FAILED: [
            "The analysis will continue without color information",
            "Consider using standard theme colors",
            "Export and re-import the presentation"
        ],
        ErrorType.API_ERROR: [
            "Check your OpenAI API key",
            "Verify you have sufficient API credits",
            "Try again in a few moments"
        ],
        ErrorType.MEMORY_ERROR: [
            "Try processing a smaller file",
            "Close other applications to free up memory",
            "Process fewer slides at once"
        ]
    }
    return suggestions.get(error_type, ["Please try again or contact support"])

# Flask error handlers
@app.errorhandler(PresentationAnalyzerError)
def handle_custom_error(error):
    """Handle custom presentation analyzer errors"""
    response = {
        'error': error.message,
        'type': error.error_type.value,
        'details': error.details,
        'suggestions': get_error_suggestions(error.error_type)
    }
    
    status_code = get_error_status_code(error.error_type)
    return jsonify(response), status_code

# Configuration
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024 * 3  # 600MB max file size
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'your-secret-key-here')

# Ensure upload directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Allowed file extensions
ALLOWED_EXTENSIONS = {'ppt', 'pptx', 'doc', 'docx', 'pdf', 'txt'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

class FileParser:
    """Enhanced file parser with better PPT parsing capabilities"""
    
    def __init__(self):
        self.supported_formats = {
            'ppt': self.parse_powerpoint,
            'pptx': self.parse_powerpoint,
            'doc': self.parse_word,
            'docx': self.parse_word,
            'pdf': self.parse_pdf,
            'txt': self.parse_text
        }
    
    def parse_file(self, file_path, file_type):
        """Parse file based on its type"""
        if file_type not in self.supported_formats:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return self.supported_formats[file_type](file_path)
    
    def parse_powerpoint(self, file_path):
        """Enhanced PowerPoint parsing with detailed slide analysis and design elements"""
        try:
            prs = Presentation(file_path)
            
            slides_data = []
            total_text = ""
            slide_count = len(prs.slides)
            
            # Track design elements across all slides
            all_colors = set()
            all_fonts = set()
            all_layouts = set()
            bullet_points = 0
            total_images = 0
            total_charts = 0
            total_tables = 0
            
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                slide_notes = ""
                shapes_count = len(slide.shapes)
                text_boxes = 0
                images = 0
                charts = 0
                tables = 0
                slide_colors = set()
                slide_fonts = set()
                
                # Extract slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide_notes = slide.notes_slide.notes_text_frame.text
                
                # Process shapes (text boxes, images, etc.)
                for shape in slide.shapes:
                    # Try multiple ways to extract text from shapes
                    shape_text = ""
                    
                    # Method 1: Direct text attribute
                    if hasattr(shape, 'text') and shape.text.strip():
                        shape_text = shape.text.strip()
                    
                    # Method 2: Text frame approach (more reliable)
                    elif hasattr(shape, 'text_frame') and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                shape_text += paragraph.text.strip() + "\n"
                    
                    # Method 3: Try to get text from any text-like shape
                    elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        if shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    shape_text += paragraph.text.strip() + "\n"
                    
                    # If we found text, process it
                    if shape_text.strip():
                        slide_text += shape_text + "\n"
                        text_boxes += 1
                        
                        # Count bullet points
                        if '•' in shape_text or '·' in shape_text or '- ' in shape_text:
                            bullet_points += shape_text.count('•') + shape_text.count('·') + shape_text.count('- ')
                        
                        # Extract font information if available
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                if hasattr(paragraph, 'runs') and paragraph.runs:
                                    for run in paragraph.runs:
                                        if hasattr(run, 'font') and run.font.name:
                                            slide_fonts.add(run.font.name)
                                            all_fonts.add(run.font.name)
                                        
                                        # Extract color information with better error handling
                                        if hasattr(run, 'font') and run.font.color:
                                            try:
                                                if hasattr(run.font.color, 'rgb'):
                                                    color = run.font.color.rgb
                                                    if color:
                                                        slide_colors.add(f"#{color:06X}")
                                                        all_colors.add(f"#{color:06X}")
                                                elif hasattr(run.font.color, 'theme_color'):
                                                    # Handle theme colors
                                                    theme_color = run.font.color.theme_color
                                                    if theme_color:
                                                        slide_colors.add(f"theme_{theme_color}")
                                                        all_colors.add(f"theme_{theme_color}")
                                            except Exception as color_error:
                                                # Skip color extraction if it fails
                                                logger.debug(f"Color extraction failed: {color_error}")
                                                pass
                    
                    # Count other shape types
                    elif shape.shape_type == 13:  # Picture
                        images += 1
                        total_images += 1
                    elif shape.shape_type == 17:  # Chart
                        charts += 1
                        total_charts += 1
                    elif shape.shape_type == 19:  # Table
                        tables += 1
                        total_tables += 1
                
                # Get layout information
                layout_type = self._get_layout_type(slide)
                all_layouts.add(layout_type)
                
                slide_data = {
                    'slide_number': i + 1,
                    'text_content': slide_text.strip(),
                    'notes': slide_notes.strip(),
                    'shapes_count': shapes_count,
                    'text_boxes': text_boxes,
                    'images': images,
                    'charts': charts,
                    'tables': tables,
                    'layout_type': layout_type,
                    'colors': list(slide_colors),
                    'fonts': list(slide_fonts)
                }
                
                slides_data.append(slide_data)
                total_text += slide_text + "\n\n"
            
            # Extract presentation metadata
            metadata = self._extract_presentation_metadata(prs)
            
            # Add design analysis metadata
            metadata.update({
                'colors': list(all_colors),
                'fonts': list(all_fonts),
                'layouts': list(all_layouts),
                'bullet_point_count': bullet_points,
                'total_images': total_images,
                'total_charts': total_charts,
                'total_tables': total_tables,
                'has_images': total_images > 0,
                'has_charts': total_charts > 0,
                'has_tables': total_tables > 0,
                'design_complexity': self._assess_design_complexity(slides_data),
                'content_density': self._assess_content_density(total_text, slide_count)
            })
            
            return {
                'slide_count': slide_count,
                'text_content': total_text.strip(),
                'slides': slides_data,
                'metadata': metadata,
                'parse_method': 'python-pptx enhanced parsing with design analysis'
            }
            
        except Exception as e:
            logger.error(f"Error parsing PowerPoint file: {str(e)}")
            raise ValueError(f"PowerPoint parsing failed: {str(e)}")
    
    def _assess_design_complexity(self, slides_data):
        """Assess the overall design complexity of the presentation"""
        total_shapes = sum(slide['shapes_count'] for slide in slides_data)
        total_images = sum(slide['images'] for slide in slides_data)
        total_charts = sum(slide['charts'] for slide in slides_data)
        total_tables = sum(slide['tables'] for slide in slides_data)
        
        complexity_score = 0
        if total_shapes > 50:
            complexity_score += 3
        elif total_shapes > 20:
            complexity_score += 2
        else:
            complexity_score += 1
            
        if total_images > 10:
            complexity_score += 2
        elif total_images > 5:
            complexity_score += 1
            
        if total_charts > 5:
            complexity_score += 2
        elif total_charts > 0:
            complexity_score += 1
            
        if total_tables > 3:
            complexity_score += 2
        elif total_tables > 0:
            complexity_score += 1
            
        if complexity_score >= 6:
            return "high"
        elif complexity_score >= 3:
            return "medium"
        else:
            return "low"
    
    def _assess_content_density(self, text_content, slide_count):
        """Assess the content density of the presentation"""
        words_per_slide = len(text_content.split()) / slide_count if slide_count > 0 else 0
        
        if words_per_slide > 100:
            return "high"
        elif words_per_slide > 50:
            return "medium"
        else:
            return "low"
    
    def _get_layout_type(self, slide):
        """Determine slide layout type"""
        try:
            layout_name = slide.slide_layout.name
            if 'title' in layout_name.lower():
                return 'title'
            elif 'content' in layout_name.lower():
                return 'content'
            elif 'section' in layout_name.lower():
                return 'section'
            else:
                return 'custom'
        except:
            return 'unknown'
    
    def _extract_presentation_metadata(self, prs):
        """Extract presentation metadata"""
        metadata = {
            'title': '',
            'author': '',
            'subject': '',
            'keywords': '',
            'category': '',
            'comments': '',
            'total_slides': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height
        }
        
        try:
            core_props = prs.core_properties
            if core_props.title:
                metadata['title'] = core_props.title
            if core_props.author:
                metadata['author'] = core_props.author
            if core_props.subject:
                metadata['subject'] = core_props.subject
            if core_props.keywords:
                metadata['keywords'] = core_props.keywords
            if core_props.category:
                metadata['category'] = core_props.category
            if core_props.comments:
                metadata['comments'] = core_props.comments
        except:
            pass
        
        return metadata
    
    def parse_word(self, file_path):
        """Parse Word documents"""
        try:
            with open(file_path, 'rb') as docx_file:
                result = mammoth.extractRawText(docx_file)
                text_content = result.value
                
                # Extract additional metadata
                metadata = {
                    'word_count': len(text_content.split()),
                    'paragraph_count': text_content.count('\n\n') + 1,
                    'has_images': 'image' in text_content.lower(),
                    'parse_method': 'mammoth (Word)'
                }
                
                return {
                    'slide_count': None,
                    'text_content': text_content,
                    'slides': [],
                    'metadata': metadata
                }
        except Exception as e:
            logger.error(f"Error parsing Word file: {str(e)}")
            raise ValueError(f"Word document parsing failed: {str(e)}")
    
    def parse_pdf(self, file_path):
        """Parse PDF files"""
        try:
            text_content = ""
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    text_content += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                
                metadata = {
                    'page_count': len(pdf_reader.pages),
                    'word_count': len(text_content.split()),
                    'parse_method': 'PyPDF2'
                }
                
                return {
                    'slide_count': len(pdf_reader.pages),
                    'text_content': text_content,
                    'slides': [],
                    'metadata': metadata
                }
        except Exception as e:
            logger.error(f"Error parsing PDF file: {str(e)}")
            raise ValueError(f"PDF parsing failed: {str(e)}")
    
    def parse_text(self, file_path):
        """Parse text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as text_file:
                text_content = text_file.read()
                
                metadata = {
                    'word_count': len(text_content.split()),
                    'line_count': len(text_content.splitlines()),
                    'parse_method': 'direct text'
                }
                
                return {
                    'slide_count': None,
                    'text_content': text_content,
                    'slides': [],
                    'metadata': metadata
                }
        except Exception as e:
            logger.error(f"Error parsing text file: {str(e)}")
            raise ValueError(f"Text file parsing failed: {str(e)}")


class EnhancedFileParser:
    """Enhanced PowerPoint parser with chunking and parallel processing"""
    
    def __init__(self, chunk_size: int = 10):
        self.chunk_size = chunk_size
        self.executor = ThreadPoolExecutor(max_workers=4)
        self.skip_formatting = False # New flag to skip formatting extraction
        
    def parse_powerpoint_chunked(self, file_path: str) -> Dict[str, Any]:
        """Parse large PowerPoint files with chunking and streaming"""
        try:
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            # Process metadata first
            metadata = self._extract_presentation_metadata_safe(prs)
            
            # Process slides in chunks
            slides_data = []
            for chunk_start in range(0, total_slides, self.chunk_size):
                chunk_end = min(chunk_start + self.chunk_size, total_slides)
                chunk_data = self._process_slide_chunk(
                    prs, chunk_start, chunk_end
                )
                slides_data.extend(chunk_data)
                
                # Yield progress for streaming (optional)
                progress = (chunk_end / total_slides) * 100
                logger.info(f"Processing progress: {progress:.1f}%")
            
            # Aggregate results
            return self._aggregate_results(slides_data, metadata, total_slides)
            
        except Exception as e:
            logger.error(f"Error parsing PowerPoint: {str(e)}")
            raise ValueError(f"PowerPoint parsing failed: {str(e)}")
    
    def _process_slide_chunk(self, prs, start_idx: int, end_idx: int) -> List[Dict]:
        """Process a chunk of slides"""
        chunk_results = []
        
        for i in range(start_idx, end_idx):
            slide = prs.slides[i]
            slide_data = self._process_single_slide_safe(slide, i + 1)
            chunk_results.append(slide_data)
            
        return chunk_results
    
    def _process_single_slide_safe(self, slide, slide_number: int) -> Dict:
        """Safely process a single slide with error handling"""
        slide_data = {
            'slide_number': slide_number,
            'text_content': '',
            'notes': '',
            'shapes_count': 0,
            'text_boxes': 0,
            'images': 0,
            'charts': 0,
            'tables': 0,
            'layout_type': 'unknown',
            'colors': [],
            'fonts': []
        }
        
        try:
            # Extract text content safely
            slide_text = []
            slide_colors = set()
            slide_fonts = set()
            
            # Get slide notes safely
            try:
                if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                    if hasattr(slide.notes_slide, 'notes_text_frame'):
                        slide_data['notes'] = slide.notes_slide.notes_text_frame.text
            except Exception as e:
                logger.debug(f"Notes extraction failed: {e}")
            
            # Process shapes with timeout protection
            slide_data['shapes_count'] = len(slide.shapes) if hasattr(slide, 'shapes') else 0
            
            for shape in slide.shapes:
                try:
                    # Extract text safely
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        if shape.text_frame:
                            text = self._extract_text_from_shape(shape)
                            if text:
                                slide_text.append(text)
                                slide_data['text_boxes'] += 1
                    
                    # Count shape types safely
                    if hasattr(shape, 'shape_type'):
                        shape_type = shape.shape_type
                        if shape_type == 13:  # Picture
                            slide_data['images'] += 1
                        elif shape_type == 17:  # Chart
                            slide_data['charts'] += 1
                        elif shape_type == 19:  # Table
                            slide_data['tables'] += 1
                    
                    # Extract colors and fonts with better error handling
                    self._extract_shape_formatting_safe(shape, slide_colors, slide_fonts)
                    
                except Exception as shape_error:
                    logger.debug(f"Shape processing error: {shape_error}")
                    continue
            
            slide_data['text_content'] = '\n'.join(slide_text)
            slide_data['colors'] = list(slide_colors)[:10]  # Limit colors
            slide_data['fonts'] = list(slide_fonts)[:5]  # Limit fonts
            
            # Get layout type safely
            try:
                if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name'):
                    slide_data['layout_type'] = slide.slide_layout.name
            except:
                pass
                
        except Exception as e:
            logger.error(f"Error processing slide {slide_number}: {e}")
            
        return slide_data
    
    def _extract_text_from_shape(self, shape) -> str:
        """Safely extract text from a shape"""
        text_parts = []
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if hasattr(paragraph, 'text'):
                        text_parts.append(paragraph.text)
        except Exception as e:
            logger.debug(f"Text extraction error: {e}")
        
        return '\n'.join(text_parts).strip()
    
    def _extract_shape_formatting_safe(self, shape, colors: set, fonts: set):
        """Safely extract formatting information"""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return
                
            for paragraph in shape.text_frame.paragraphs:
                if not hasattr(paragraph, 'runs'):
                    continue
                    
                for run in paragraph.runs:
                    # Extract font safely
                    try:
                        if hasattr(run, 'font') and hasattr(run.font, 'name'):
                            if run.font.name:
                                fonts.add(run.font.name)
                    except:
                        pass
                    
                    # Extract color safely - simplified approach
                    try:
                        if hasattr(run, 'font') and hasattr(run.font, 'color'):
                            color_obj = run.font.color
                            
                            # Try RGB color
                            if hasattr(color_obj, 'rgb') and color_obj.rgb:
                                try:
                                    # Handle RGBColor object
                                    if hasattr(color_obj.rgb, '__iter__'):
                                        # It's already an RGB tuple
                                        r, g, b = color_obj.rgb
                                        colors.add(f"#{r:02X}{g:02X}{b:02X}")
                                    else:
                                        # It might be an integer
                                        rgb_int = int(color_obj.rgb)
                                        colors.add(f"#{rgb_int:06X}")
                                except (ValueError, TypeError, AttributeError):
                                    pass
                            
                            # Try theme color
                            elif hasattr(color_obj, 'theme_color') and color_obj.theme_color:
                                colors.add(f"theme_{color_obj.theme_color}")
                    except:
                        pass
                        
        except Exception as e:
            logger.debug(f"Formatting extraction error: {e}")
    
    def _extract_presentation_metadata_safe(self, prs) -> Dict:
        """Safely extract presentation metadata"""
        metadata = {
            'title': '',
            'author': '',
            'subject': '',
            'total_slides': len(prs.slides),
            'slide_width': prs.slide_width if hasattr(prs, 'slide_width') else None,
            'slide_height': prs.slide_height if hasattr(prs, 'slide_height') else None
        }
        
        try:
            if hasattr(prs, 'core_properties'):
                core_props = prs.core_properties
                metadata['title'] = getattr(core_props, 'title', '') or ''
                metadata['author'] = getattr(core_props, 'author', '') or ''
                metadata['subject'] = getattr(core_props, 'subject', '') or ''
        except Exception as e:
            logger.debug(f"Metadata extraction error: {e}")
            
        return metadata
    
    def _aggregate_results(self, slides_data: List[Dict], metadata: Dict, total_slides: int) -> Dict:
        """Aggregate slide data into final result"""
        # Collect all text
        all_text = []
        all_colors = set()
        all_fonts = set()
        total_images = 0
        total_charts = 0
        total_tables = 0
        bullet_points = 0
        
        for slide in slides_data:
            if slide['text_content']:
                all_text.append(slide['text_content'])
                # Count bullet points
                bullet_points += slide['text_content'].count('•')
                bullet_points += slide['text_content'].count('·')
                bullet_points += slide['text_content'].count('- ')
            
            all_colors.update(slide.get('colors', []))
            all_fonts.update(slide.get('fonts', []))
            total_images += slide.get('images', 0)
            total_charts += slide.get('charts', 0)
            total_tables += slide.get('tables', 0)
        
        # Update metadata
        metadata.update({
            'colors': list(all_colors)[:20],  # Limit to 20 colors
            'fonts': list(all_fonts)[:10],  # Limit to 10 fonts
            'bullet_point_count': bullet_points,
            'total_images': total_images,
            'total_charts': total_charts,
            'total_tables': total_tables,
            'has_images': total_images > 0,
            'has_charts': total_charts > 0,
            'has_tables': total_tables > 0,
            'design_complexity': self._assess_design_complexity(slides_data),
            'content_density': self._assess_content_density('\n'.join(all_text), total_slides)
        })
        
        return {
            'slide_count': total_slides,
            'text_content': '\n\n'.join(all_text),
            'slides': slides_data,
            'metadata': metadata,
            'parse_method': 'Enhanced chunked parsing with error handling'
        }
    
    def _assess_design_complexity(self, slides_data: List[Dict]) -> str:
        """Assess design complexity based on slide elements"""
        total_elements = sum(
            slide.get('shapes_count', 0) for slide in slides_data
        )
        avg_elements = total_elements / len(slides_data) if slides_data else 0
        
        if avg_elements > 15:
            return 'high'
        elif avg_elements > 8:
            return 'medium'
        else:
            return 'low'
    
    def _assess_content_density(self, text: str, slide_count: int) -> str:
        """Assess content density"""
        if slide_count == 0:
            return 'low'
            
        words_per_slide = len(text.split()) / slide_count
        
        if words_per_slide > 100:
            return 'high'
        elif words_per_slide > 50:
            return 'medium'
        else:
            return 'low'


class OpenAIService:
    """OpenAI service for design analysis"""
    
    def __init__(self, api_key):
        self.api_key = api_key
        openai.api_key = api_key
    
    def analyze_design(self, content, company_context="", audience_info=None):
        """Analyze design using OpenAI GPT-4"""
        try:
            prompt = self._create_design_compass_prompt(content, company_context, audience_info)
            
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[
                    {
                        "role": "system",
                        "content": """You are a senior presentation design consultant following "The Presentation Design Compass" methodology. You help designers move from strategy to design through a 5-stage framework:

1. Understand the context (contextual grounding)
2. Set the ambitions (clarify expectations) 
3. Map out decisions (visual direction planning)
4. Tell the story (execute with storytelling principles)
5. Sell your work (justify and present)

You classify presentations into 4 types:
- Executive & Strategic: Formal, concise, outcome-driven for leadership
- Sales & Influence: Bold, compelling, benefit-led for persuasion
- Engagement & Immersion: Uplifting, emotional, visual for inspiration
- Informative & Educational: Clear, helpful, structured for explanation

Always respond with valid JSON only. Consider presentation type when making design recommendations."""
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                max_tokens=3000,
                temperature=0.7
            )
            
            analysis_text = response.choices[0].message.content
            
            # Clean and parse JSON response
            cleaned_response = analysis_text.replace('```json\n', '').replace('```\n', '').strip()
            
            return json.loads(cleaned_response)
            
        except Exception as e:
            logger.error(f"OpenAI analysis error: {str(e)}")
            raise ValueError(f"Analysis failed: {str(e)}")
    
    def analyze_slide(self, slide_data, presentation_context, audience_info, slide_number):
        """Analyze individual slide with detailed recommendations"""
        try:
            prompt = self._create_slide_analysis_prompt(slide_data, presentation_context, audience_info, slide_number)
            
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[
                    {
                        "role": "system",
                        "content": """You are an expert presentation design consultant specializing in slide-by-slide analysis. You provide detailed, actionable recommendations for individual slides that help designers create compelling, effective presentations.

Your expertise includes:
- Slide layout optimization and visual hierarchy
- Content organization and readability
- Color theory and typography for individual slides
- Visual storytelling and narrative flow
- Audience engagement and retention
- Technical design principles and best practices

IMPORTANT: Always respond with valid JSON format. Do not include any text outside the JSON structure."""
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                max_tokens=2000,
                temperature=0.7
            )
            
            analysis_text = response.choices[0].message.content
            
            # Clean and parse JSON response with better error handling
            try:
                # Remove markdown code blocks if present
                cleaned_response = analysis_text.replace('```json\n', '').replace('```\n', '').replace('```json', '').replace('```', '').strip()
                
                # Try to parse the JSON
                return json.loads(cleaned_response)
                
            except json.JSONDecodeError as json_error:
                logger.error(f"JSON parsing error: {json_error}")
                logger.error(f"Raw response: {analysis_text}")
                
                # Return a fallback analysis structure
                return {
                    "slideOverview": {
                        "slideNumber": slide_number,
                        "contentSummary": "Analysis failed - please try again",
                        "slidePurpose": "Unable to determine",
                        "effectiveness": "unknown",
                        "priority": "unknown"
                    },
                    "contentAnalysis": {
                        "textContent": {
                            "clarity": "Unable to analyze",
                            "organization": "Unable to analyze",
                            "length": "unknown",
                            "keyMessages": ["Analysis failed"],
                            "improvements": ["Please try analyzing this slide again"]
                        },
                        "visualElements": {
                            "images": {
                                "count": slide_data.get('images', 0),
                                "relevance": "Unable to analyze",
                                "quality": "unknown",
                                "recommendations": ["Please try again"]
                            },
                            "charts": {
                                "count": slide_data.get('charts', 0),
                                "effectiveness": "Unable to analyze",
                                "clarity": "Unable to analyze",
                                "improvements": ["Please try again"]
                            },
                            "tables": {
                                "count": slide_data.get('tables', 0),
                                "readability": "Unable to analyze",
                                "structure": "Unable to analyze",
                                "improvements": ["Please try again"]
                            }
                        }
                    },
                    "designRecommendations": {
                        "layout": {
                            "currentLayout": slide_data.get('layout_type', 'unknown'),
                            "effectiveness": "Unable to analyze",
                            "recommendedLayout": "Unable to determine",
                            "specificChanges": ["Please try analyzing this slide again"],
                            "visualHierarchy": "Unable to analyze"
                        },
                        "colorScheme": {
                            "currentColors": slide_data.get('colors', []),
                            "recommendations": ["Unable to analyze"],
                            "contrast": "Unable to analyze",
                            "accessibility": "Unable to analyze"
                        },
                        "typography": {
                            "currentFonts": slide_data.get('fonts', []),
                            "readability": "Unable to analyze",
                            "recommendations": ["Unable to analyze"],
                            "hierarchy": "Unable to analyze"
                        }
                    },
                    "actionItems": {
                        "immediate": ["Try analyzing this slide again"],
                        "shortTerm": ["Contact support if issue persists"],
                        "longTerm": ["Consider manual review"]
                    },
                    "error": "JSON parsing failed - please try again"
                }
            
        except Exception as e:
            logger.error(f"OpenAI slide analysis error: {str(e)}")
            raise ValueError(f"Slide analysis failed: {str(e)}")
    
    def _create_design_compass_prompt(self, content, company_context, audience_info=None):
        """Create the design compass analysis prompt"""
        
        # Build content sections
        slide_count_text = ""
        if content.get('slide_count'):
            slide_count_text = f"Slide Count: {content['slide_count']}"
        
        colors_text = ""
        if content.get('metadata', {}).get('colors'):
            colors_text = f"Current Colors: {', '.join(content['metadata']['colors'])}"
        
        fonts_text = ""
        if content.get('metadata', {}).get('fonts'):
            fonts_text = f"Fonts: {', '.join(content['metadata']['fonts'])}"
        
        images_text = ""
        if content.get('metadata', {}).get('has_images'):
            images_text = f"Has Images: {content['metadata']['has_images']}"
        
        bullets_text = ""
        if content.get('metadata', {}).get('bullet_point_count'):
            bullets_text = f"Bullet Points: {content['metadata']['bullet_point_count']}"
        
        company_text = ""
        if company_context:
            company_text = f"COMPANY DESIGN GUIDELINES:\n{company_context}\n"
        
        audience_text = ""
        if audience_info:
            retainer_client = audience_info.get('retainer_client', '')
            retainer_info = f"- Retainer Client: {retainer_client}" if retainer_client else ""
            
            audience_text = f"""AUDIENCE INFORMATION:
- Type: {audience_info.get('type', 'Not specified')}
- Goal: {audience_info.get('goal', 'Not specified')}
- Size: {audience_info.get('size', 'Not specified')}
- Additional Context: {audience_info.get('context', 'None provided')}
{retainer_info}

"""
        
        return f"""Following "The Presentation Design Compass" methodology, analyze this presentation comprehensively:

DOCUMENT CONTENT:
{slide_count_text}
Text Content: {content.get('text_content', '')}
{colors_text}
{fonts_text}
{images_text}
{bullets_text}

{company_text}
{audience_text}

Provide a comprehensive design analysis with specific strategic direction, color palettes, and design elements. Respond with this detailed JSON structure:

{{
  "presentationType": {{
    "primary": "Executive & Strategic | Sales & Influence | Engagement & Immersion | Informative & Educational",
    "secondary": "optional secondary type or null",
    "reasoning": "detailed explanation of why this classification based on content analysis",
    "confidence": "high | medium | low"
  }},
  "strategicDirection": {{
    "primaryStrategy": "specific strategic approach (e.g., 'Data-driven executive summary for board approval')",
    "communicationGoal": "what the presentation should achieve",
    "audienceEngagement": "how to engage the specific audience",
    "callToAction": "what action should the audience take",
    "successMetrics": ["how to measure presentation success"]
  }},
  "contextualGrounding": {{
    "identifiedObjective": "persuade | inform | align | inspire",
    "audienceProfile": "detailed description of likely audience including their role, background, and expectations",
    "presentationTrigger": "why this presentation exists now - the business context",
    "stakeholderMapping": "key stakeholders and their specific interests/concerns",
    "urgencyLevel": "high | medium | low",
    "politicalClimate": "any sensitivities or dynamics detected",
    "businessContext": "industry, market conditions, and organizational context"
  }},
  "designDirection": {{
    "backgrounds": {{
      "recommended": "White/Light | Dark | Colored | Textured/Thematic | Gradient | Minimal",
      "specificRecommendations": ["detailed background suggestions"],
      "reasoning": "why this background approach fits the presentation type and audience",
      "examples": ["specific background styles to consider"]
    }},
    "layouts": {{
      "recommended": "Structured & repetitive | Dynamic & alternating | Asymmetrical | Full-screen visuals | Grid-based | Free-form",
      "specificRecommendations": ["detailed layout suggestions"],
      "reasoning": "how this supports the content and audience",
      "slideStructure": ["specific slide layout recommendations"]
    }},
    "imagery": {{
      "recommended": "Photography | Illustration | 3D/isometric | Hybrid | None | Icons | Charts",
      "specificRecommendations": ["detailed imagery suggestions"],
      "reasoning": "why this imagery approach works for this type",
      "imageStyle": "specific image style recommendations",
      "chartTypes": ["recommended chart types if data is present"]
    }},
    "fonts": {{
      "headings": "specific font recommendation with reasoning",
      "body": "specific font recommendation with reasoning", 
      "accent": "accent font for highlights",
      "reasoning": "how typography supports presentation type and tone",
      "fontPairings": ["recommended font combinations"],
      "sizing": "font size recommendations"
    }},
    "colors": {{
      "primary": ["#specific color codes with names"],
      "secondary": ["#specific color codes with names"],
      "accent": ["#accent color codes"],
      "neutral": ["#neutral color codes"],
      "paletteApproach": "Blues (trust) | Greens (growth) | Reds (energy) | Black/dark (premium) | Pastels (calm) | Gradients (innovation) | Corporate | Brand-specific",
      "colorPsychology": "detailed explanation of color choices and their psychological impact",
      "accessibility": "color contrast and accessibility considerations",
      "brandAlignment": "how colors align with brand guidelines"
    }},
    "visualMetaphors": {{
      "recommended": "specific metaphor suggestions based on content",
      "transformationTheme": "what's changing from → to",
      "reasoning": "how metaphors support the story",
      "visualElements": ["specific visual elements to include"]
    }},
    "spacing": {{
      "recommended": "tight | moderate | generous",
      "reasoning": "spacing strategy explanation",
      "specificGuidelines": ["spacing recommendations"]
    }},
    "visualHierarchy": {{
      "recommended": "clear hierarchy strategy",
      "reasoning": "how to create visual hierarchy",
      "specificGuidelines": ["hierarchy recommendations"]
    }}
  }},
  "storytellingStructure": {{
    "narrativeApproach": "Problem-Solution | Journey | Comparison | Transformation | Data Story | Hero's Journey | Before-After-Bridge",
    "keyMessages": ["main message 1", "main message 2", "main message 3"],
    "emotionalTone": "confident | urgent | aspirational | calm | disruptive | authoritative | inspiring",
    "flowRecommendations": ["specific slide flow suggestions"],
    "opening": "how to start the presentation",
    "closing": "how to end the presentation",
    "transitions": ["transition recommendations between sections"]
  }},
  "contentStrategy": {{
    "keyPoints": ["main content points to emphasize"],
    "dataPresentation": "how to present any data or statistics",
    "storyElements": ["storytelling elements to include"],
    "callToAction": "specific call to action recommendations",
    "supportingEvidence": "what evidence or proof points to include"
  }},
  "executionGuidance": {{
    "priorityFixes": ["top 3 priority improvements in order with specific details"],
    "quickWins": ["easy improvements with high impact"],
    "designPrinciples": ["key principles to follow for this presentation type"],
    "slideTemplateNeeds": ["what slide templates would help most"],
    "technicalSpecs": ["technical specifications and requirements"],
    "deliveryTips": ["presentation delivery recommendations"]
  }},
  "clientQuestions": {{
    "clarifyingQuestions": ["questions to ask client about context/ambitions"],
    "stakeholderQuestions": ["questions about audience and dynamics"],
    "visualReadinessQuestions": ["questions about design preferences and constraints"],
    "technicalQuestions": ["questions about technical requirements"],
    "timelineQuestions": ["questions about project timeline and milestones"]
  }},
  "designCompassStage": {{
    "currentStage": "Context | Ambitions | Decisions | Story | Sell",
    "nextSteps": ["what to do next in the design process"],
    "stageGuidance": "specific advice for moving to next stage",
    "deliverables": ["what deliverables are needed for this stage"]
  }},
  "brandIntegration": {{
    "brandAlignment": "how to align with brand guidelines",
    "brandElements": ["specific brand elements to include"],
    "brandColors": ["brand color recommendations"],
    "brandVoice": "tone and voice recommendations"
  }},
    "technicalSpecifications": {{
    "format": "PowerPoint | Keynote | PDF | Web",
    "aspectRatio": "16:9 | 4:3 | custom",
    "resolution": "recommended resolution",
    "fileSize": "target file size considerations",
    "compatibility": "compatibility requirements"
  }}
}}"""

    def _create_slide_analysis_prompt(self, slide_data, presentation_context, audience_info, slide_number):
        """Create the slide-specific analysis prompt"""
        
        # Build slide content sections
        slide_text = slide_data.get('text_content', '')
        slide_notes = slide_data.get('notes', '')
        shapes_count = slide_data.get('shapes_count', 0)
        text_boxes = slide_data.get('text_boxes', 0)
        images = slide_data.get('images', 0)
        charts = slide_data.get('charts', 0)
        tables = slide_data.get('tables', 0)
        layout_type = slide_data.get('layout_type', 'unknown')
        colors = slide_data.get('colors', [])
        fonts = slide_data.get('fonts', [])
        
        # Build context sections
        audience_text = ""
        if audience_info:
            retainer_client = audience_info.get('retainer_client', '')
            retainer_info = f"- Retainer Client: {retainer_client}" if retainer_client else ""
            
            audience_text = f"""AUDIENCE INFORMATION:
- Type: {audience_info.get('type', 'Not specified')}
- Goal: {audience_info.get('goal', 'Not specified')}
- Size: {audience_info.get('size', 'Not specified')}
- Additional Context: {audience_info.get('context', 'None provided')}
{retainer_info}

"""
        
        presentation_type = presentation_context.get('presentationType', {}).get('primary', 'Unknown')
        strategic_goal = presentation_context.get('strategicDirection', {}).get('primaryStrategy', 'Not specified')
        
        return f"""Analyze this individual slide (Slide {slide_number}) with detailed, actionable recommendations:

SLIDE CONTENT:
Text Content: {slide_text}
Speaker Notes: {slide_notes}
Layout Type: {layout_type}
Shapes Count: {shapes_count}
Text Boxes: {text_boxes}
Images: {images}
Charts: {charts}
Tables: {tables}
Colors Used: {', '.join(colors) if colors else 'None detected'}
Fonts Used: {', '.join(fonts) if fonts else 'Default fonts'}

PRESENTATION CONTEXT:
Presentation Type: {presentation_type}
Strategic Goal: {strategic_goal}

{audience_text}

Provide detailed, slide-specific analysis and recommendations. Respond with this JSON structure:

{{
  "slideOverview": {{
    "slideNumber": {slide_number},
    "contentSummary": "brief summary of what this slide contains",
    "slidePurpose": "what this slide is trying to achieve",
    "effectiveness": "high | medium | low",
    "priority": "critical | important | nice-to-have"
  }},
  "contentAnalysis": {{
    "textContent": {{
      "clarity": "how clear and readable the text is",
      "organization": "how well the content is structured",
      "length": "appropriate | too long | too short",
      "keyMessages": ["main points this slide conveys"],
      "improvements": ["specific text improvements"]
    }},
    "visualElements": {{
      "images": {{
        "count": {images},
        "relevance": "how well images support the content",
        "quality": "high | medium | low",
        "recommendations": ["image-specific suggestions"]
      }},
      "charts": {{
        "count": {charts},
        "effectiveness": "how well data is presented",
        "clarity": "how easy to understand",
        "improvements": ["chart-specific suggestions"]
      }},
      "tables": {{
        "count": {tables},
        "readability": "how easy to scan",
        "structure": "how well organized",
        "improvements": ["table-specific suggestions"]
      }}
    }}
  }},
  "designRecommendations": {{
    "layout": {{
      "currentLayout": "{layout_type}",
      "effectiveness": "how well the layout works",
      "recommendedLayout": "suggested layout type",
      "specificChanges": ["layout improvement suggestions"],
      "visualHierarchy": "how to improve visual flow"
    }},
    "colorScheme": {{
      "currentColors": {colors},
      "effectiveness": "how well colors work together",
      "recommendedPalette": ["specific color suggestions"],
      "colorPsychology": "why these colors work for this audience",
      "accessibility": "color contrast considerations"
    }},
    "typography": {{
      "currentFonts": {fonts},
      "readability": "how easy to read",
      "recommendedFonts": ["font suggestions"],
      "sizing": "font size recommendations",
      "hierarchy": "how to create better text hierarchy"
    }},
    "spacing": {{
      "currentSpacing": "tight | moderate | generous",
      "recommendations": ["spacing improvements"],
      "whiteSpace": "how to use white space better"
    }}
  }},
  "storytelling": {{
    "narrativeRole": "how this slide fits the story",
    "flow": "how it connects to previous/next slides",
    "engagement": "how to make it more engaging",
    "callToAction": "what action this slide should inspire"
  }},
  "actionItems": {{
    "immediate": ["quick fixes that can be done right away"],
    "shortTerm": ["improvements for the next iteration"],
    "longTerm": ["strategic changes for future versions"]
  }},
  "priority": {{
    "level": "critical | high | medium | low",
    "reasoning": "why this slide needs attention",
    "impact": "what improving this slide will achieve"
  }}
}}"""

# Initialize services
file_parser = FileParser()

@app.route('/')
def index():
    """Main application page - serve React app"""
    try:
        # Try to serve from build directory (production)
        return send_from_directory('build', 'index.html')
    except FileNotFoundError:
        # Fallback to public directory (development)
        return send_file('public/index.html')

@app.route('/<path:path>')
def serve_static(path):
    """Serve static files from build directory"""
    try:
        return send_from_directory('build', path)
    except FileNotFoundError:
        # Fallback to public directory
        return send_from_directory('public', path)

@app.route('/health')
def health_check():
    """Health check endpoint"""
    return jsonify({
        'status': 'healthy', 
        'timestamp': datetime.now().isoformat(),
        'services': {
            'openai': 'available',
            'file_parser': 'available',
            'memory_usage': MemoryManager.check_memory_usage()
        }
    })

@app.route('/parse-pptx', methods=['POST'])
def parse_pptx():
    """Enhanced PowerPoint parsing endpoint with comprehensive error handling"""
    try:
        # Check rate limiting
        if not rate_limiter.can_make_call():
            return jsonify({
                'error': 'Rate limit exceeded. Please wait before making another request.',
                'type': 'RATE_LIMIT_EXCEEDED'
            }), 429
        
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Validate file before processing
        validation_result = FileValidator.validate_file(file)
        if not validation_result['valid']:
            return jsonify({
                'error': 'File validation failed',
                'details': validation_result['errors'],
                'warnings': validation_result['warnings']
            }), 400
        
        # Log warnings if any
        if validation_result['warnings']:
            for warning in validation_result['warnings']:
                logger.warning(f"File warning: {warning['message']}")
        
        # Save file temporarily
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Check file size for logging
        file_size = validation_result['file_info']['size']
        if file_size > 100 * 1024 * 1024: # > 100MB
            logger.info(f"Large file detected: {file_size / (1024*1024):.2f} MB")
        
        try:
            # Use memory management decorator
            @MemoryManager.process_with_memory_management
            def process_file():
                file_extension = validation_result['file_info']['extension']
                if file_extension in ['ppt', 'pptx']:
                    parser = EnhancedFileParser(chunk_size=10)
                    return parser.parse_powerpoint_chunked(file_path)
                else:
                    return file_parser.parse_file(file_path, file_extension)
            
            result = process_file()
            
            response_size = len(json.dumps(result))
            if response_size > 10 * 1024 * 1024: # > 10MB response
                logger.warning(f"Large response detected: {response_size / (1024*1024):.2f} MB")
                summary = {
                    'slide_count': result['slide_count'],
                    'text_content': result['text_content'][:5000] + '...' if len(result['text_content']) > 5000 else result['text_content'],
                    'metadata': result['metadata'],
                    'slides_available': True,
                    'total_slides': len(result.get('slides', [])),
                    'response_truncated': True,
                    'original_size_mb': response_size / (1024*1024)
                }
                return jsonify(summary)
            
            logger.info(f"Parse result - slide_count: {result.get('slide_count')}, text_length: {len(result.get('text_content', ''))}")
            return jsonify(result)
        finally:
            if os.path.exists(file_path):
                os.remove(file_path)
    except PresentationAnalyzerError as e:
        logger.error(f"Presentation analyzer error: {e.message}")
        return jsonify({
            'error': e.message,
            'type': e.error_type.value,
            'details': e.details,
            'suggestions': get_error_suggestions(e.error_type)
        }), get_error_status_code(e.error_type)
    except Exception as e:
        logger.error(f"Unexpected error in parse_pptx: {str(e)}")
        return jsonify({'error': 'An unexpected error occurred while processing the file'}), 500

@app.route('/analyze', methods=['POST'])
def analyze_design():
    """Analyze design using OpenAI"""
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No data provided'}), 400
        
        api_key = data.get('api_key')
        if not api_key:
            return jsonify({'error': 'OpenAI API key required'}), 400
        
        content = data.get('content')
        company_context = data.get('company_context', '')
        audience_info = data.get('audience_info', {})
        slides_data = data.get('slides_data', [])  # Add slides data
        
        if not content:
            return jsonify({'error': 'Content required for analysis'}), 400
        
        # Initialize OpenAI service
        openai_service = OpenAIService(api_key)
        
        # Analyze design with audience information
        analysis = openai_service.analyze_design(content, company_context, audience_info)
        
        # Include slides data in the response for slide-by-slide analysis
        analysis['slides'] = slides_data
        
        return jsonify(analysis)
        
    except Exception as e:
        logger.error(f"Error in analyze_design: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/export/pdf', methods=['POST'])
def export_pdf():
    """Export analysis to PDF"""
    try:
        data = request.get_json()
        analysis = data.get('analysis')
        
        if not analysis:
            return jsonify({'error': 'Analysis data required'}), 400
        
        # Create PDF content (simplified for now)
        pdf_content = _create_pdf_content(analysis)
        
        # For now, return as text file
        return send_file(
            io.BytesIO(pdf_content.encode()),
            mimetype='text/plain',
            as_attachment=True,
            download_name='design-analysis-report.txt'
        )
        
    except Exception as e:
        logger.error(f"Error in export_pdf: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/chat', methods=['POST'])
def chat():
    """Chat with the design assistant"""
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No data provided'}), 400
        
        api_key = data.get('api_key')
        if not api_key:
            return jsonify({'error': 'OpenAI API key required'}), 400
        
        message = data.get('message')
        analysis = data.get('analysis')
        context = data.get('context', '')
        conversation_history = data.get('conversation_history', [])
        
        if not message:
            return jsonify({'error': 'Message required'}), 400
        
        # Initialize OpenAI service
        openai_service = OpenAIService(api_key)
        
        # Create enhanced chat prompt with better context management
        system_prompt = """You are an expert presentation design consultant with deep knowledge of "The Presentation Design Compass" methodology. You help designers create compelling presentations that achieve their strategic goals.

Your expertise includes:
- Strategic presentation planning and audience analysis
- Color theory and visual design principles
- Typography and layout optimization
- Storytelling and narrative structure
- Brand integration and visual hierarchy
- Technical specifications and delivery considerations

Always provide specific, actionable advice based on the analysis context. Be encouraging but direct, and focus on practical improvements that will have the most impact."""

        # Build context summary instead of dumping full JSON
        context_summary = ""
        if analysis:
            context_summary = f"""PRESENTATION ANALYSIS SUMMARY:
- Type: {analysis.get('presentationType', {}).get('primary', 'Unknown')}
- Strategic Goal: {analysis.get('strategicDirection', {}).get('primaryStrategy', 'Not specified')}
- Audience: {analysis.get('contextualGrounding', {}).get('audienceProfile', 'Not specified')}
- Key Design Elements: {', '.join(analysis.get('designDirection', {}).get('colors', {}).get('primary', [])[:3]) if analysis.get('designDirection', {}).get('colors', {}).get('primary') else 'Not specified'}
- Priority Issues: {', '.join(analysis.get('executionGuidance', {}).get('priorityFixes', [])[:2]) if analysis.get('executionGuidance', {}).get('priorityFixes') else 'Not specified'}
"""

        company_context_summary = ""
        if context:
            company_context_summary = f"COMPANY CONTEXT: {context}"

        # Build conversation history for context
        conversation_context = ""
        if conversation_history:
            recent_messages = conversation_history[-3:]  # Last 3 messages for context
            conversation_context = f"RECENT CONVERSATION:\n" + "\n".join([
                f"{'User' if msg.get('role') == 'user' else 'Assistant'}: {msg.get('content', '')}"
                for msg in recent_messages
            ])

        # Create the user prompt
        user_prompt = f"""{context_summary}
{company_context_summary}
{conversation_context}

USER QUESTION: {message}

Please provide a helpful, specific response about their presentation design. Focus on practical advice and actionable insights. If they're asking about specific aspects of the analysis, reference the relevant data. If they need clarification on any design principles, explain them clearly."""
        
        # Prepare messages for OpenAI
        messages = [{"role": "system", "content": system_prompt}]
        
        # Add conversation history (limit to last 10 messages to avoid token limits)
        if conversation_history:
            for msg in conversation_history[-10:]:
                messages.append({
                    "role": msg.get('role', 'user'),
                    "content": msg.get('content', '')
                })
        
        # Add current message
        messages.append({"role": "user", "content": user_prompt})
        
        # Get response from OpenAI
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=1000,
            temperature=0.7
        )
        
        chat_response = response.choices[0].message.content
        
        return jsonify({'response': chat_response})
        
    except Exception as e:
        logger.error(f"Error in chat: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/export/json', methods=['POST'])
def export_json():
    """Export analysis to JSON"""
    try:
        data = request.get_json()
        analysis = data.get('analysis')
        
        if not analysis:
            return jsonify({'error': 'Analysis data required'}), 400
        
        return send_file(
            io.BytesIO(json.dumps(analysis, indent=2).encode()),
            mimetype='application/json',
            as_attachment=True,
            download_name='design-analysis.json'
        )
        
    except Exception as e:
        logger.error(f"Error in export_json: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/analyze-slide', methods=['POST'])
def analyze_slide():
    """Analyze individual slide with detailed recommendations"""
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No data provided'}), 400
        
        api_key = data.get('api_key')
        if not api_key:
            return jsonify({'error': 'OpenAI API key required'}), 400
        
        slide_data = data.get('slide_data')
        presentation_context = data.get('presentation_context', {})
        audience_info = data.get('audience_info', {})
        slide_number = data.get('slide_number', 1)
        
        if not slide_data:
            return jsonify({'error': 'Slide data required for analysis'}), 400
        
        # Initialize OpenAI service
        openai_service = OpenAIService(api_key)
        
        # Analyze individual slide
        slide_analysis = openai_service.analyze_slide(slide_data, presentation_context, audience_info, slide_number)
        
        return jsonify(slide_analysis)
        
    except Exception as e:
        logger.error(f"Error in analyze_slide: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/export/questions', methods=['POST'])
def export_questions():
    """Export client questions"""
    try:
        data = request.get_json()
        analysis = data.get('analysis')
        
        if not analysis:
            return jsonify({'error': 'Analysis data required'}), 400
        
        questions_content = f"""Client Questions for Design Discussion

Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

CLARIFYING QUESTIONS:
{chr(10).join([f"- {q}" for q in analysis.get('clientQuestions', {}).get('clarifyingQuestions', [])])}

STAKEHOLDER QUESTIONS:
{chr(10).join([f"- {q}" for q in analysis.get('clientQuestions', {}).get('stakeholderQuestions', [])])}

VISUAL READINESS QUESTIONS:
{chr(10).join([f"- {q}" for q in analysis.get('clientQuestions', {}).get('visualReadinessQuestions', [])])}
"""
        
        return send_file(
            io.BytesIO(questions_content.encode()),
            mimetype='text/plain',
            as_attachment=True,
            download_name='client-questions.txt'
        )
        
    except Exception as e:
        logger.error(f"Error in export_questions: {str(e)}")
        return jsonify({'error': str(e)}), 500

def _create_pdf_content(analysis):
    """Create PDF content from analysis"""
    content = f"""
AI Design Analysis Report
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

PRESENTATION TYPE
Primary: {analysis.get('presentationType', {}).get('primary', 'N/A')}
Secondary: {analysis.get('presentationType', {}).get('secondary', 'N/A')}
Reasoning: {analysis.get('presentationType', {}).get('reasoning', 'N/A')}

CONTEXTUAL GROUNDING
Objective: {analysis.get('contextualGrounding', {}).get('identifiedObjective', 'N/A')}
Audience: {analysis.get('contextualGrounding', {}).get('audienceProfile', 'N/A')}
Urgency: {analysis.get('contextualGrounding', {}).get('urgencyLevel', 'N/A')}

DESIGN DIRECTION
Backgrounds: {analysis.get('designDirection', {}).get('backgrounds', {}).get('recommended', 'N/A')}
Layouts: {analysis.get('designDirection', {}).get('layouts', {}).get('recommended', 'N/A')}
Imagery: {analysis.get('designDirection', {}).get('imagery', {}).get('recommended', 'N/A')}

STORYTELLING STRUCTURE
Narrative: {analysis.get('storytellingStructure', {}).get('narrativeApproach', 'N/A')}
Tone: {analysis.get('storytellingStructure', {}).get('emotionalTone', 'N/A')}

EXECUTION GUIDANCE
Priority Fixes:
{chr(10).join([f"- {fix}" for fix in analysis.get('executionGuidance', {}).get('priorityFixes', [])])}

Quick Wins:
{chr(10).join([f"- {win}" for win in analysis.get('executionGuidance', {}).get('quickWins', [])])}

CLIENT QUESTIONS
Clarifying Questions:
{chr(10).join([f"- {q}" for q in analysis.get('clientQuestions', {}).get('clarifyingQuestions', [])])}

Stakeholder Questions:
{chr(10).join([f"- {q}" for q in analysis.get('clientQuestions', {}).get('stakeholderQuestions', [])])}
"""
    return content

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000) 